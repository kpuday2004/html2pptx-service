#!/usr/bin/env python3

import http.server
import io
import os
import base64
from bs4 import BeautifulSoup
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.oxml.ns import qn
from lxml import etree

SLIDE_W = 13.33
SLIDE_H = 7.5
MARGIN = 0.5
server_port = int(os.getenv('HTML2PPTX_PORT', '8080'))

# ── Color palette (Light Theme) ────────────────────────────────────────────
BG       = RGBColor(0xF8, 0xFA, 0xFC)   # near white background
BG2      = RGBColor(0xF1, 0xF5, 0xF9)   # light gray card background
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
ACCENT   = RGBColor(0x25, 0x63, 0xEB)   # strong blue
ACCENT2  = RGBColor(0x05, 0x96, 0x69)   # emerald green
ACCENT3  = RGBColor(0xD9, 0x77, 0x06)   # amber
ACCENT4  = RGBColor(0xDB, 0x27, 0x77)   # pink
ACCENT5  = RGBColor(0x70, 0x3B, 0xE5)   # purple
SUBTEXT  = RGBColor(0x64, 0x74, 0x8B)   # slate gray
BODYTEXT = RGBColor(0x1E, 0x29, 0x3B)   # dark navy text

CHART_COLORS = [ACCENT, ACCENT2, ACCENT3, ACCENT4, ACCENT5]


# ── SHARED HELPERS ──────────────────────────────────────────────────────────

def set_bg(slide, color=None):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color or BG


def add_shape(slide, x, y, w, h, color):
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_textbox(slide, text, x, y, w, h, size=16, bold=False, color=WHITE,
                italic=False, align=None, wrap=True, auto_size=True):
    from pptx.enum.text import PP_ALIGN
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    if auto_size:
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.italic = italic
    p.font.color.rgb = color
    if align:
        p.alignment = align
    return tb


def add_accent_bar(slide, x=None, y=1.45, w=None, h=0.04):
    x = x if x is not None else MARGIN
    w = w if w is not None else (SLIDE_W - 2 * MARGIN)
    bar = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background()


def add_title(slide, text, y=0.4, size=36, color=BODYTEXT):
    add_textbox(slide, text, MARGIN, y, SLIDE_W - 2 * MARGIN, 1.0,
                size=size, bold=True, color=color)


def add_bullets(slide, items, x=None, y=1.6, w=None, size=16):
    x = x if x is not None else MARGIN
    w = w if w is not None else (SLIDE_W - 2 * MARGIN)
    tb = slide.shapes.add_textbox(
        Inches(x), Inches(y), Inches(w),
        Inches(SLIDE_H - y - MARGIN)
    )
    tf = tb.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = '  •  ' + item
        p.font.size = Pt(size)
        p.font.color.rgb = BODYTEXT
        p.space_before = Pt(6)


def add_paragraph_text(slide, text, x=None, y=1.6, w=None, h=0.8):
    x = x if x is not None else MARGIN
    w = w if w is not None else (SLIDE_W - 2 * MARGIN)
    add_textbox(slide, text, x, y, w, h, size=15, color=SUBTEXT, italic=True)


def apply_table_borders(tbl):
    """Apply light borders to all cells in a table."""
    for row in tbl.rows:
        for cell in row.cells:
            tc = cell._tc
            tcPr = tc.get_or_add_tcPr()
            for border_tag in ['a:lnL', 'a:lnR', 'a:lnT', 'a:lnB']:
                existing = tcPr.find(qn(border_tag))
                if existing is not None:
                    tcPr.remove(existing)
                border_el = etree.fromstring(
                    f'<{border_tag} xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" w="6350">'
                    f'<a:solidFill><a:srgbClr val="E2E8F0"/></a:solidFill>'
                    f'</{border_tag}>'
                )
                tcPr.append(border_el)


def add_chart_table(slide, categories, values, series_name, x, y, w, h):
    """Adds a compact data table to accompany a chart."""
    from pptx.enum.text import PP_ALIGN

    categories = list(categories)[:10]
    values = list(values)[:10]
    row_count = len(categories) + 1  # +1 for header
    col_count = 2

    tbl = slide.shapes.add_table(
        row_count, col_count,
        Inches(x), Inches(y),
        Inches(w), Inches(h)
    ).table

    col_w = Inches(w / col_count)
    for col in tbl.columns:
        col.width = col_w

    # Header row
    headers = [series_name or 'Category', 'Value']
    for j, header in enumerate(headers):
        cell = tbl.cell(0, j)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = ACCENT
        para = cell.text_frame.paragraphs[0]
        para.font.bold = True
        para.font.size = Pt(11)
        para.font.color.rgb = WHITE
        para.alignment = PP_ALIGN.CENTER

    # Data rows
    for i, (cat, val) in enumerate(zip(categories, values)):
        row_color = WHITE if i % 2 == 0 else BG2

        cat_cell = tbl.cell(i + 1, 0)
        cat_cell.text = str(cat)
        cat_cell.fill.solid()
        cat_cell.fill.fore_color.rgb = row_color
        cat_para = cat_cell.text_frame.paragraphs[0]
        cat_para.font.size = Pt(10)
        cat_para.font.color.rgb = BODYTEXT
        cat_para.alignment = PP_ALIGN.LEFT

        val_cell = tbl.cell(i + 1, 1)
        val_str = str(int(val)) if val == int(val) else f'{val:.1f}'
        val_cell.text = val_str
        val_cell.fill.solid()
        val_cell.fill.fore_color.rgb = row_color
        val_para = val_cell.text_frame.paragraphs[0]
        val_para.font.size = Pt(10)
        val_para.font.color.rgb = ACCENT
        val_para.font.bold = True
        val_para.alignment = PP_ALIGN.CENTER

    apply_table_borders(tbl)


def build_chart(slide, chart_type_enum, chart_title, categories, series_name, values, x, y, w, h):
    """Build a styled chart and inject an XML legend."""
    chart_data = ChartData()
    chart_data.categories = categories
    chart_data.add_series(series_name, values)

    chart = slide.shapes.add_chart(
        chart_type_enum,
        Inches(x), Inches(y), Inches(w), Inches(h),
        chart_data
    ).chart

    chart.has_title = True
    chart.chart_title.text_frame.text = chart_title
    chart.chart_title.text_frame.paragraphs[0].font.color.rgb = BODYTEXT
    chart.chart_title.text_frame.paragraphs[0].font.size = Pt(13)

    plot = chart.plots[0]
    plot.has_data_labels = True
    plot.data_labels.font.color.rgb = BODYTEXT
    plot.data_labels.font.size = Pt(10)

    if chart_type_enum == XL_CHART_TYPE.PIE:
        plot.data_labels.number_format = '0"%"'

    for i, point in enumerate(plot.series[0].points):
        point.format.fill.solid()
        point.format.fill.fore_color.rgb = CHART_COLORS[i % len(CHART_COLORS)]

    try:
        chart.category_axis.tick_labels.font.color.rgb = SUBTEXT
        chart.value_axis.tick_labels.font.color.rgb = SUBTEXT
        chart.category_axis.format.line.color.rgb = SUBTEXT
        chart.value_axis.format.line.color.rgb = SUBTEXT
    except Exception:
        pass

    # Inject legend via XML
    chart_xml = chart._element
    plot_area = chart_xml.find(qn('c:plotArea'))
    if plot_area is not None:
        legend_xml = etree.fromstring('''
            <c:legend xmlns:c="http://schemas.openxmlformats.org/drawingml/2006/chart">
                <c:legendPos val="b"/>
                <c:overlay val="0"/>
                <c:spPr>
                    <a:noFill xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"/>
                </c:spPr>
                <c:txPr xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">
                    <a:bodyPr/>
                    <a:lstStyle/>
                    <a:p>
                        <a:pPr>
                            <a:defRPr sz="1000" b="0">
                                <a:solidFill>
                                    <a:srgbClr val="64748B"/>
                                </a:solidFill>
                            </a:defRPr>
                        </a:pPr>
                    </a:p>
                </c:txPr>
            </c:legend>
        ''')
        chart_xml.insert(list(chart_xml).index(plot_area) + 1, legend_xml)

    return chart


# ── SLIDE BUILDERS ──────────────────────────────────────────────────────────

def build_title_slide(prs, el):
    """LAYOUT 1 — Title slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)

    add_shape(slide, 0, 0, SLIDE_W, 0.08, ACCENT)
    add_shape(slide, 0, SLIDE_H - 1.2, SLIDE_W, 1.2, RGBColor(0xE2, 0xE8, 0xF0))
    add_shape(slide, 0, SLIDE_H - 1.22, SLIDE_W, 0.05, ACCENT)
    add_shape(slide, 0, 0.08, 0.06, SLIDE_H - 1.28, ACCENT2)

    title_el = el.find('h1')
    title_text = title_el.get_text().strip() if title_el else 'Presentation'

    tb = slide.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(SLIDE_W - 1.0), Inches(2.5))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0x1E, 0x29, 0x3B)

    sub_el = el.find('p')
    if sub_el:
        add_textbox(slide, sub_el.get_text().strip(),
                    0.5, 4.5, SLIDE_W - 1.0, 0.7,
                    size=20, color=ACCENT, italic=True)

    from pptx.enum.text import PP_ALIGN
    add_textbox(slide, 'Powered by Infinium Intelligence',
                0.5, SLIDE_H - 1.0, SLIDE_W - 1.0, 0.5,
                size=11, color=BODYTEXT, align=PP_ALIGN.RIGHT)


def build_bullets_slide(prs, el):
    """LAYOUT 2 — Standard bullets slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_shape(slide, 0, 0, SLIDE_W, 0.05, ACCENT)

    title_el = el.find('h1')
    add_title(slide, title_el.get_text().strip() if title_el else 'Slide')
    add_accent_bar(slide)

    items = [li.get_text().strip() for li in el.find_all('li')]
    if items:
        add_bullets(slide, items)


def build_bar_chart_slide(prs, el):
    """LAYOUT 3 — Bar chart with inline data table."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_shape(slide, 0, 0, SLIDE_W, 0.05, ACCENT)

    title_el = el.find('h1')
    add_title(slide, title_el.get_text().strip() if title_el else 'Chart')
    add_accent_bar(slide)

    cats = [c.strip() for c in el.get('data-categories', '').split(',') if c.strip()]
    vals = [float(v.strip()) for v in el.get('data-values', '0').split(',') if v.strip()]
    series_name = el.get('data-series-name', '')
    chart_title = el.get('data-chart-title', '')

    content_y = 1.55
    p_el = el.find('p')
    if p_el:
        add_paragraph_text(slide, p_el.get_text().strip(), y=content_y, h=0.55)
        content_y = 2.2

    content_h = SLIDE_H - content_y - MARGIN

    if len(cats) <= 6:
        # Chart on top (~65%), table below (~35%)
        chart_h = content_h * 0.62
        table_y = content_y + chart_h + 0.1
        table_h = content_h - chart_h - 0.1

        build_chart(slide, XL_CHART_TYPE.COLUMN_CLUSTERED,
                    chart_title, cats, series_name, vals,
                    MARGIN, content_y, SLIDE_W - 2 * MARGIN, chart_h)

        add_chart_table(slide, cats, vals, series_name,
                        MARGIN + 2.0, table_y,
                        SLIDE_W - 2 * MARGIN - 4.0, table_h)
    else:
        # Chart left (~65%), table right (~35%)
        chart_w = (SLIDE_W - 2 * MARGIN) * 0.65
        table_x = MARGIN + chart_w + 0.2
        table_w = SLIDE_W - table_x - MARGIN

        build_chart(slide, XL_CHART_TYPE.COLUMN_CLUSTERED,
                    chart_title, cats, series_name, vals,
                    MARGIN, content_y, chart_w, content_h)

        add_chart_table(slide, cats, vals, series_name,
                        table_x, content_y, table_w, content_h)


def build_pie_chart_slide(prs, el):
    """LAYOUT 4 — Pie chart with inline data table."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_shape(slide, 0, 0, SLIDE_W, 0.05, ACCENT)

    title_el = el.find('h1')
    add_title(slide, title_el.get_text().strip() if title_el else 'Chart')
    add_accent_bar(slide)

    cats = [c.strip() for c in el.get('data-categories', '').split(',') if c.strip()]
    vals = [float(v.strip()) for v in el.get('data-values', '0').split(',') if v.strip()]
    chart_title = el.get('data-chart-title', '')

    content_y = 1.55
    p_el = el.find('p')
    if p_el:
        add_paragraph_text(slide, p_el.get_text().strip(), y=content_y, h=0.55)
        content_y = 2.2

    content_h = SLIDE_H - content_y - MARGIN

    # Pie always: chart left (~62%), table right (~38%)
    chart_w = (SLIDE_W - 2 * MARGIN) * 0.62
    table_x = MARGIN + chart_w + 0.2
    table_w = SLIDE_W - table_x - MARGIN

    build_chart(slide, XL_CHART_TYPE.PIE,
                chart_title, cats, '', vals,
                MARGIN, content_y, chart_w, content_h)

    add_chart_table(slide, cats, vals, chart_title,
                    table_x, content_y, table_w, content_h)


def build_two_column_slide(prs, el):
    """LAYOUT 5 — Two column comparison slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_shape(slide, 0, 0, SLIDE_W, 0.05, ACCENT)

    title_el = el.find('h1')
    add_title(slide, title_el.get_text().strip() if title_el else 'Comparison')
    add_accent_bar(slide)

    col_w = (SLIDE_W - 2 * MARGIN - 0.3) / 2
    left = el.find(class_='col-left')
    right = el.find(class_='col-right')

    for col_x, col_el, accent_color in [
        (MARGIN, left, ACCENT),
        (MARGIN + col_w + 0.3, right, ACCENT2)
    ]:
        if not col_el:
            continue
        add_shape(slide, col_x, 1.65, col_w, SLIDE_H - 1.65 - MARGIN, BG2)
        add_shape(slide, col_x, 1.65, col_w, 0.05, accent_color)

        heading_el = col_el.find('h2')
        if heading_el:
            add_textbox(slide, heading_el.get_text().strip(),
                        col_x + 0.2, 1.78, col_w - 0.4, 0.55,
                        size=17, bold=True, color=accent_color)

        items = [li.get_text().strip() for li in col_el.find_all('li')]
        if items:
            add_bullets(slide, items, x=col_x + 0.2, y=2.45,
                        w=col_w - 0.4, size=14)


def build_mixed_slide(prs, el):
    """LAYOUT 6 — Bullets left, chart top-right, table bottom-right."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_shape(slide, 0, 0, SLIDE_W, 0.05, ACCENT)

    title_el = el.find('h1')
    add_title(slide, title_el.get_text().strip() if title_el else 'Overview')
    add_accent_bar(slide)

    content_y = 1.6
    content_h = SLIDE_H - content_y - MARGIN
    col_w = (SLIDE_W - 2 * MARGIN - 0.4) / 2

    # Left: bullet card
    add_shape(slide, MARGIN, content_y, col_w, content_h, BG2)
    add_shape(slide, MARGIN, content_y, col_w, 0.04, ACCENT2)

    col_left = el.find(class_='col-left')
    items = [li.get_text().strip() for li in (col_left or el).find_all('li')]
    if items:
        add_bullets(slide, items, x=MARGIN + 0.2, y=content_y + 0.2,
                    w=col_w - 0.4, size=14)

    # Right: chart top 65%, table bottom 35%
    chart_x = MARGIN + col_w + 0.4
    chart_w = SLIDE_W - chart_x - MARGIN
    chart_h = content_h * 0.62
    table_y = content_y + chart_h + 0.1
    table_h = content_h - chart_h - 0.1

    chart_type_str = el.get('data-chart', 'bar')
    chart_type_enum = XL_CHART_TYPE.PIE if chart_type_str == 'pie' else XL_CHART_TYPE.COLUMN_CLUSTERED

    cats = [c.strip() for c in el.get('data-categories', '').split(',') if c.strip()]
    vals_raw = [v.strip() for v in el.get('data-values', '').split(',') if v.strip()]
    vals = [float(v) for v in vals_raw] if vals_raw else [1]
    series_name = el.get('data-series-name', '')
    chart_title = el.get('data-chart-title', '')

    if cats and vals:
        build_chart(slide, chart_type_enum,
                    chart_title, cats, series_name, vals,
                    chart_x, content_y, chart_w, chart_h)

        add_chart_table(slide, cats, vals, series_name,
                        chart_x, table_y, chart_w, table_h)


# ── MAIN CONVERTER ──────────────────────────────────────────────────────────

def html_to_pptx_bytes(html_string):
    soup = BeautifulSoup(html_string, 'html.parser')
    slide_els = soup.select('section.slide')

    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)

    for el in slide_els:
        chart_type = el.get('data-chart', '').strip()
        layout = el.get('data-layout', '').strip()

        is_title = (
            not chart_type and
            not layout and
            el.find('p') and
            not el.find('ul')
        )

        if is_title:
            build_title_slide(prs, el)
        elif chart_type == 'bar' and layout != 'mixed':
            build_bar_chart_slide(prs, el)
        elif chart_type == 'pie' and layout != 'mixed':
            build_pie_chart_slide(prs, el)
        elif layout == 'two-column':
            build_two_column_slide(prs, el)
        elif layout == 'mixed':
            build_mixed_slide(prs, el)
        else:
            build_bullets_slide(prs, el)

    if not slide_els:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        set_bg(slide)
        add_textbox(slide, 'No slides detected — check HTML structure',
                    1, 3, 11, 1, size=20, color=RGBColor(0xFF, 0x55, 0x55))

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


class Handler(http.server.BaseHTTPRequestHandler):
    def do_POST(self):
        content_length = int(self.headers.get('Content-Length', 0))
        body = self.rfile.read(content_length).decode('utf-8').strip()

        try:
            html_string = base64.b64decode(body).decode('utf-8')
            print("Decoded as base64")
        except Exception:
            html_string = body
            print("Used as raw HTML")

        print(f"HTML length: {len(html_string)}")
        print(f"First 100 chars: {html_string[:100]}")

        pptx_bytes = html_to_pptx_bytes(html_string)

        self.send_response(200)
        self.send_header('Content-Type', 'application/vnd.openxmlformats-officedocument.presentationml.presentation')
        self.send_header('Content-Disposition', 'attachment; filename="presentation.pptx"')
        self.send_header('Content-Length', str(len(pptx_bytes)))
        self.send_header('Access-Control-Allow-Origin', '*')
        self.end_headers()
        self.wfile.write(pptx_bytes)

    def do_OPTIONS(self):
        self.send_response(200)
        self.send_header('Access-Control-Allow-Origin', '*')
        self.send_header('Access-Control-Allow-Methods', 'POST, OPTIONS')
        self.send_header('Access-Control-Allow-Headers', 'Content-Type')
        self.end_headers()

    def log_message(self, format, *args):
        pass


httpd = http.server.HTTPServer(('', server_port), Handler)
print(f'html2pptx serving on port {server_port}')
httpd.serve_forever()
